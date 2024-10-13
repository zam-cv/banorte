import PyPDF2
from docx import Document
from pptx import Presentation
import re
import pytesseract
import cv2

def archivo_procesado(archivo:str)->str:
    # Función para identificar el tipo de archivo
    def identificar_tipo_archivo(archivo_binario):
        if archivo_binario.startswith(b'%PDF'):
            return 'pdf'
        elif archivo_binario.startswith(b'PK'):
            # DOCX y PPTX son archivos ZIP que comienzan con 'PK'
            if b'word/' in archivo_binario:
                return 'docx'
            elif b'ppt/' in archivo_binario:
                return 'pptx'
        elif archivo_binario.startswith((b'\xff\xd8', b'\x89PNG')):
            return 'image'
        else:
            return 'txt'

    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

    with open(archivo, 'rb') as file:
        archivo_binario = file.read()
        
    tipo_archivo = identificar_tipo_archivo(archivo_binario)
    print(tipo_archivo)


    regex = re.compile(r'\t|\n|\r|  ')
    # Abrir archivo dependiendo del tipo de archivo
    texto:list[str]=[]
    if(tipo_archivo == 'txt'):
        fichero = open(archivo, 'r')
        texto = fichero.readlines()
    elif(tipo_archivo == 'pdf'):
        # Abre el archivo PDF en modo lectura binaria
        with open(archivo, 'rb') as fichero:
            # Crea un objeto de lectura de PDF
            lector_pdf = PyPDF2.PdfReader(fichero)
            # Obtén el número de páginas en el PDF
            num_paginas = len(lector_pdf.pages)
            # Lee cada página y extrae el texto
            for i in range(num_paginas):
                pagina = lector_pdf.pages[i]
                text = pagina.extract_text().split('\n')
                texto=text
    elif(tipo_archivo == 'docx'):
        # Abre el archivo DOCX
        documento = Document(archivo)
        for parrafo in documento.paragraphs:
            texto.append(parrafo.text)
    elif(tipo_archivo == 'pptx'):
        # Abre el archivo de PowerPoint
        presentacion = Presentation(archivo)
        # Lee y extrae el texto de cada diapositiva
        for i, diapositiva in enumerate(presentacion.slides):
            for forma in diapositiva.shapes:
                if hasattr(forma, "text"):
                    texto.append(forma.text)
    elif(tipo_archivo == 'image'):
        # Leer la imagen
        img = cv2.imread('a.jpg')

        # Convertir la imagen a escala de grises
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

        # Aplicar umbral binario
        _, binary = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY_INV)

        # Aplicar dilatación y erosión para eliminar el ruido
        kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, 1))
        dilated = cv2.dilate(binary, kernel, iterations=1)
        eroded = cv2.erode(dilated, kernel, iterations=1)

        # Extraer texto de la imagen usando Tesseract OCR
        texto = pytesseract.image_to_string(eroded, lang='spa')  # Cambia 'spa' por el idioma que necesites
        
    texto_filtrado=[]
    if tipo_archivo != 'image':
        for linea in texto:
            nueva_linea = ''.join(char if not regex.match(char) else ' ' for char in linea)
            prev_char = ''
            nueva_linea2 = ''
            for char in nueva_linea:
                if char == ' ' and prev_char == ' ':
                    continue
                else:
                    nueva_linea2+=(char)
                prev_char = char
            texto_filtrado.append(nueva_linea2)
        texto_final=''
        for linea in texto_filtrado:
            texto_final+=linea+'' if linea.endswith(' ') else linea+' '
    else:
        texto_final=""
        for char in texto:
            if not regex.match(char):
                texto_final+=char
        
    print(texto_final)
    return texto_final

if __name__ == '__main__':
    #archivo_procesado('a.jpg')
    #archivo_procesado('ejemplo.txt')
    #archivo_procesado('prueba.pdf')
    #archivo_procesado('prueba.docx')
    archivo_procesado('ejemplo.pptx')